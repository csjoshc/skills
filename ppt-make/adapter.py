#!/usr/bin/env python3
"""
ppt-make Adapter v2 — Rich Content Engine.

Converts YAML slide definitions into C3-branded PPTX decks with:
- Flow diagrams (boxes + directional connectors + branching)
- Architecture diagrams (layered boxes with layer labels + legend)
- Pipeline diagrams (horizontal strip + detail tables)
- Panel columns (2-4 columns with colored headers)
- Comparison layouts (side-by-side diagrams)
- Tables with auto column widths and font scaling
- Auto font sizing based on content density
- In-slide hyperlinks and speaker notes

Usage:
    python adapter.py input.yaml [--template path/to/template.pptx] [--output deck.pptx]

    # Also accepts markdown with YAML frontmatter (legacy mode):
    python adapter.py input.md
"""
from __future__ import annotations

import argparse
import math
import os
import re
import sys
from typing import Any

from pptx import Presentation
from pptx.util import Inches as I, Pt as P
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# YAML / frontmatter parsing
# ---------------------------------------------------------------------------
try:
    import yaml
    _HAS_YAML = True
except ImportError:
    _HAS_YAML = False


def _load_yaml(text: str) -> dict[str, Any]:
    if _HAS_YAML:
        return yaml.safe_load(text) or {}
    raise ImportError("PyYAML is required for the v2 adapter. pip install pyyaml")


def _parse_input(path: str) -> dict[str, Any]:
    """Parse a .yaml or .md file into a deck spec."""
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()

    if path.endswith((".yaml", ".yml")):
        return _load_yaml(text)

    # Legacy: markdown with YAML frontmatter
    fm = re.match(r"^---\s*\n(.*?)\n---\s*\n", text, re.DOTALL)
    if fm:
        meta = _load_yaml(fm.group(1))
        body = text[fm.end():]
    else:
        meta = {}
        body = text

    # Convert ## Slide: headers to slides list
    parts = re.split(r"(?m)^##\s+Slide:\s*(.+)", body)
    slides = []
    for i in range(1, len(parts), 2):
        title = parts[i].strip()
        content = parts[i + 1].strip() if i + 1 < len(parts) else ""
        slides.append({"title": title, "content": content})

    meta["slides"] = slides
    return meta


# ═══════════════════════════════════════════════════════════════
# Color palette — imported from palette.py (single source of truth)
# ═══════════════════════════════════════════════════════════════

# Add skill dir to path so palette.py is importable
_skill_dir = os.path.dirname(os.path.abspath(__file__))
if _skill_dir not in sys.path:
    sys.path.insert(0, _skill_dir)

from palette import (  # noqa: E402
    COLORS as _COLORS,
    color as _color,
    WHITE, LIGHT_GRAY, DARK_FILL, C3_BLUE, ORANGE, GREEN,
    WARM_GRAY, PANEL_DARK, GRAY_FILL, TBL_HDR, TBL_ROW_A, TBL_ROW_B,
)


# ═══════════════════════════════════════════════════════════════
# Template loader
# ═══════════════════════════════════════════════════════════════

def _load_template(path: str):
    """Load template, fix footer year, return (prs, layout_fn)."""
    import datetime
    prs = Presentation(path)
    year = str(datetime.date.today().year)
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "2025" in run.text:
                            run.text = run.text.replace("2025", year)
        for lay in master.slide_layouts:
            for shape in lay.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "2025" in run.text:
                                run.text = run.text.replace("2025", year)

    def layout(name: str):
        for L in prs.slide_layouts:
            if L.name == name:
                return L
        return prs.slide_layouts[13]  # Blank fallback

    return prs, layout


# ═══════════════════════════════════════════════════════════════
# Primitive helpers
# ═══════════════════════════════════════════════════════════════

def _add_title(s, text: str, sz: int = 24):
    """White bold title at top of slide."""
    tb = s.shapes.add_textbox(I(0.30), I(0.18), I(12.75), I(0.58))
    tf = tb.text_frame
    tf.word_wrap = False
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = P(sz)
    r.font.bold = True
    r.font.color.rgb = WHITE


def _add_subtitle(s, text: str, top: float = 0.82, sz: int = 12):
    """Italic light-gray subtitle line."""
    _txt(s, text, 0.30, top, 12.0, 0.30, sz=sz, color=LIGHT_GRAY, italic=True)


def _txt(s, text: str, left: float, top: float, w: float, h: float,
         sz: int = 12, bold: bool = False, color=None,
         align=PP_ALIGN.LEFT, wrap: bool = True, italic: bool = False):
    """Add a textbox with multi-line support."""
    color = color or WHITE
    tb = s.shapes.add_textbox(I(left), I(top), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = P(sz)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def _box(s, text: str, left: float, top: float, w: float, h: float,
         fill=None, fsz: int = 12, fbold: bool = True, fcolor=None,
         border=None):
    """Colored rectangle with centered text."""
    fcolor = fcolor or WHITE
    sh = s.shapes.add_shape(1, I(left), I(top), I(w), I(h))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        sh.line.width = P(0.75)
    else:
        sh.line.fill.background()
    if text:
        tf = sh.text_frame
        tf.word_wrap = True
        for j, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = P(fsz)
            r.font.bold = fbold
            r.font.color.rgb = fcolor
    return sh


def _connector(s, x1: float, y1: float, x2: float, y2: float,
               color=None, width: float = 1.5):
    """Straight connector with arrowhead."""
    color = color or ORANGE
    conn = s.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color
    conn.line.width = P(width)
    ln = conn.line._ln
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return conn


def _c3tbl(s, left: float, top: float, w: float, h: float,
           rows: list, cw=None, fsz: int = 11, hdr_fsz: int | None = None):
    """C3-styled table with dark header and striped body."""
    hdr_fsz = hdr_fsz or fsz
    nr, nc = len(rows), len(rows[0])
    t = s.shapes.add_table(nr, nc, I(left), I(top), I(w), I(h)).table
    if cw:
        for i, c in enumerate(cw):
            t.columns[i].width = I(c)
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = t.cell(ri, ci)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else p.add_run()
            r.font.color.rgb = WHITE
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_HDR
                r.font.bold = True
                r.font.size = P(hdr_fsz)
            else:
                r.font.size = P(fsz)
                if ri % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TBL_ROW_A
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TBL_ROW_B
    return t


def _add_notes(s, text: str):
    """Set speaker notes."""
    if text:
        tf = s.notes_slide.notes_text_frame
        tf.text = str(text)


def _hlink(s, text: str, url: str, left: float, top: float,
           w: float = 3.50, h: float = 0.22, sz: int = 8, color=None):
    """Small clickable hyperlink textbox."""
    color = color or C3_BLUE
    tb = s.shapes.add_textbox(I(left), I(top), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = P(sz)
    r.font.bold = False
    r.font.color.rgb = color
    r.hyperlink.address = url
    return tb


# ═══════════════════════════════════════════════════════════════
# Auto font-sizing
# ═══════════════════════════════════════════════════════════════

def _auto_body_font(content_text: str, has_shapes: bool = False) -> int:
    """Choose body font size based on content density."""
    lines = [l for l in content_text.split("\n") if l.strip()]
    n = len(lines)
    total_chars = sum(len(l) for l in lines)
    if has_shapes or n > 10 or total_chars > 600:
        return 11
    if n > 7 or total_chars > 400:
        return 12
    if n > 4 or total_chars > 200:
        return 14
    return 16


def _auto_shape_font(num_shapes: int) -> int:
    """Choose shape label font based on shape count."""
    if num_shapes > 10:
        return 9
    if num_shapes > 7:
        return 10
    if num_shapes > 4:
        return 11
    return 12


def _auto_box_size(num_boxes: int, available_width: float = 12.40,
                   row_count: int = 1):
    """Calculate box width, height, and gap for a flow row."""
    if num_boxes <= 0:
        return 2.0, 0.85, 0.25
    # Leave room for gaps
    max_gap = 0.40
    min_gap = 0.10
    # Try to fit boxes with reasonable gaps
    gap = min(max_gap, max(min_gap, (available_width - num_boxes * 0.80) / max(num_boxes - 1, 1) * 0.3))
    bw = (available_width - (num_boxes - 1) * gap) / num_boxes
    bw = min(bw, 2.50)
    bw = max(bw, 1.00)
    gap = (available_width - num_boxes * bw) / max(num_boxes - 1, 1) if num_boxes > 1 else 0
    gap = max(gap, min_gap)
    bh = 0.90 if num_boxes <= 5 else (0.70 if num_boxes <= 8 else 0.58)
    return round(bw, 2), round(bh, 2), round(gap, 2)


# ═══════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════

def _build_title_slide(prs, layout_fn, slide_spec: dict, meta: dict):
    """Title slide using template placeholders."""
    s = prs.slides.add_slide(layout_fn("Title Slide Option"))
    title = meta.get("title", slide_spec.get("title", "Untitled"))
    subtitle = meta.get("subtitle", "")
    author = meta.get("author", "")
    date = meta.get("date", "")
    author_line = "  \u00b7  ".join(p for p in [author, date] if p)

    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
            if ph.text_frame.paragraphs[0].runs:
                ph.text_frame.paragraphs[0].runs[0].font.size = P(36)
        elif idx == 15:
            ph.text = subtitle
            if ph.text_frame.paragraphs[0].runs:
                ph.text_frame.paragraphs[0].runs[0].font.size = P(18)
        elif idx == 16:
            ph.text = author_line
            if ph.text_frame.paragraphs[0].runs:
                ph.text_frame.paragraphs[0].runs[0].font.size = P(14)

    # Sources in notes
    sources = meta.get("sources", [])
    if sources:
        lines = ["Sources:"]
        for src in sources:
            if isinstance(src, dict):
                lines.append(f"\u2022 {src.get('label','')}: {src.get('url','')}")
            else:
                lines.append(f"\u2022 {src}")
        _add_notes(s, "\n".join(lines))
    return s


def _build_panels(prs, layout_fn, slide_spec: dict, meta: dict):
    """2-4 column panel layout with colored headers."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    panels = slide_spec.get("panels", [])
    if not panels:
        return s

    num = min(len(panels), 4)
    col_gap = 0.12
    total_w = 12.40
    col_w = (total_w - (num - 1) * col_gap) / num
    panel_top = 0.92
    panel_h = 5.80
    hdr_h = 0.42

    for i, panel in enumerate(panels[:num]):
        x = 0.30 + i * (col_w + col_gap)
        hdr = panel.get("header", "")
        clr = _color(panel.get("color", "dark"))
        content = panel.get("content", "")

        # Background
        bg = s.shapes.add_shape(1, I(x), I(panel_top), I(col_w), I(panel_h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = PANEL_DARK
        bg.line.color.rgb = WARM_GRAY
        bg.line.width = P(0.5)

        # Header bar
        hb = s.shapes.add_shape(1, I(x), I(panel_top), I(col_w), I(hdr_h))
        hb.fill.solid()
        hb.fill.fore_color.rgb = clr
        hb.line.fill.background()
        _txt(s, hdr, x + 0.12, panel_top + 0.03, col_w - 0.24, 0.36,
             sz=14, bold=True)

        # Content
        body_sz = _auto_body_font(content, has_shapes=False)
        body_sz = min(body_sz, 14)  # panels cap at 14
        _txt(s, content, x + 0.15, panel_top + hdr_h + 0.16,
             col_w - 0.30, panel_h - hdr_h - 0.20, sz=body_sz)

    _add_slide_extras(s, slide_spec)
    return s


def _build_flow(prs, layout_fn, slide_spec: dict, meta: dict):
    """Flow diagram: rows of boxes with directional connectors."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    rows = slide_spec.get("rows", [])
    if not rows:
        return s

    start_y = 1.18 if not slide_spec.get("subtitle") else 1.30
    row_gap = 0.12
    available_w = 12.90

    for ri, row_spec in enumerate(rows):
        items = row_spec.get("items", [])
        if not items:
            continue

        direction = row_spec.get("direction", "right")
        bw, bh, gap = _auto_box_size(len(items), available_w)
        fsz = _auto_shape_font(len(items))
        y = start_y + ri * (bh + row_gap + 0.50)

        # Override with explicit positioning if provided
        if row_spec.get("y") is not None:
            y = row_spec["y"]

        if direction == "left":
            # Right-to-left: items[0] is rightmost
            for i, item in enumerate(items):
                x = 0.22 + (len(items) - 1 - i) * (bw + gap)
                fill = _color(item.get("fill", "dark"))
                text = item.get("text", "")
                _box(s, text, x, y, bw, bh, fill=fill, fsz=fsz)
                if i < len(items) - 1:
                    # Leftward connector
                    cx_end = x - gap
                    cx_start = x
                    _connector(s, cx_start, y + bh / 2, cx_end, y + bh / 2,
                              color=_color(item.get("connector_color", "orange")))
        else:
            for i, item in enumerate(items):
                x = 0.22 + i * (bw + gap)
                fill = _color(item.get("fill", "dark"))
                text = item.get("text", "")
                box_w = item.get("w", bw)
                _box(s, text, x, y, box_w, bh, fill=fill, fsz=fsz)
                if i < len(items) - 1:
                    cx_start = x + box_w
                    cx_end = 0.22 + (i + 1) * (bw + gap)
                    _connector(s, cx_start, y + bh / 2, cx_end, y + bh / 2,
                              color=_color(item.get("connector_color", "orange")))

        # Handle branching (YES/NO from a decision box)
        branch = row_spec.get("branch")
        if branch:
            # Find the source box index
            from_idx = branch.get("from_index", len(items) - 1)
            bx = 0.22 + from_idx * (bw + gap)
            label = branch.get("label", "YES")
            label_color = _color(branch.get("color", "orange"))
            target = branch.get("target")  # "next_row" or specific

            if branch.get("direction", "down") == "down":
                _connector(s, bx + bw / 2, y + bh,
                          bx + bw / 2, y + bh + gap + 0.08,
                          color=label_color)
                _txt(s, label, bx + bw / 2 + 0.05, y + bh + 0.02,
                     0.50, 0.18, sz=8, color=label_color, wrap=False)

            # "NO" branch to a terminal box
            no_spec = branch.get("no")
            if no_spec:
                no_text = no_spec.get("text", "COMPLETED")
                no_fill = _color(no_spec.get("fill", "green"))
                no_w = no_spec.get("w", 0.85)
                no_x = bx + bw + gap
                _box(s, no_text, no_x, y, no_w, bh, fill=no_fill, fsz=fsz)
                _connector(s, bx + bw, y + bh / 2, no_x, y + bh / 2,
                          color=no_fill)
                _txt(s, "NO", bx + bw + 0.06, y + bh / 2 - 0.22,
                     0.50, 0.18, sz=8, color=no_fill, wrap=False)

    # Output path (optional)
    output_path = slide_spec.get("output_path")
    if output_path:
        items = output_path.get("items", [])
        label = output_path.get("label", "Output path (when COMPLETED):")
        oy = output_path.get("y")
        if oy is None:
            oy = start_y + len(rows) * 1.25 + 0.40
        _txt(s, label, 0.22, oy - 0.30, 8.0, 0.28, sz=11, bold=True)
        bw, bh, gap = _auto_box_size(len(items), available_w)
        fsz = _auto_shape_font(len(items))
        for i, item in enumerate(items):
            x = 0.22 + i * (bw + gap)
            _box(s, item.get("text", ""), x, oy, bw, bh,
                 fill=_color(item.get("fill", "dark")), fsz=fsz)
            if i < len(items) - 1:
                _connector(s, x + bw, oy + bh / 2,
                          x + bw + gap, oy + bh / 2)

    # Legend
    legend = slide_spec.get("legend")
    if legend:
        _build_legend(s, legend, y_start=output_path.get("y", start_y + len(rows) * 1.25) + 1.20 if output_path else start_y + len(rows) * 1.25 + 0.50)

    _add_slide_extras(s, slide_spec)
    return s


def _build_diagram(prs, layout_fn, slide_spec: dict, meta: dict):
    """Architecture diagram with layers, boxes, edges, and legend."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    layers = slide_spec.get("layers", [])
    elements = slide_spec.get("elements", [])
    edges = slide_spec.get("edges", [])

    # Auto-layout parameters
    layer_label_w = 1.30
    content_left = 1.60
    content_w = 11.00
    layer_h = 1.10
    layer_gap = 0.40
    start_y = 1.25 if not slide_spec.get("subtitle") else 1.40

    # Build element ID → position map
    elem_map = {}
    num_layers = max(len(layers), 1)

    # Calculate layer positions
    layer_positions = {}
    for li in range(num_layers):
        ly = start_y + li * (layer_h + layer_gap)
        layer_positions[li] = ly

    # Draw layer labels (thin colored strips)
    for li, layer in enumerate(layers):
        ly = layer_positions[li]
        lbl_color = _color(layer.get("color", "dark"))
        # Thin vertical strip
        strip = s.shapes.add_shape(1, I(0.18), I(ly), I(0.14), I(layer_h))
        strip.fill.solid()
        strip.fill.fore_color.rgb = lbl_color
        strip.line.fill.background()
        _txt(s, layer.get("label", ""), 0.38, ly + 0.15, 1.10, 0.80,
             sz=10, bold=True, color=WHITE)

    # Place elements
    # Group elements by layer
    by_layer: dict[int, list] = {}
    for elem in elements:
        li = elem.get("layer", 0)
        by_layer.setdefault(li, []).append(elem)

    for li, elems in by_layer.items():
        ly = layer_positions.get(li, start_y + li * (layer_h + layer_gap))
        # Auto-layout: distribute within the layer
        n = len(elems)
        elem_gap = 0.20
        total_box_w = content_w - (n - 1) * elem_gap
        ew = total_box_w / n if n > 0 else content_w
        ew = min(ew, 5.00)  # cap width

        for ei, elem in enumerate(elems):
            # Allow explicit position overrides
            ex = elem.get("left", content_left + ei * (ew + elem_gap))
            ey = elem.get("top", ly)
            e_w = elem.get("w", ew)
            e_h = elem.get("h", layer_h - 0.20)
            fill = _color(elem.get("fill", "dark"))
            border = _color(elem.get("border")) if elem.get("border") else None
            fsz = _auto_shape_font(len(elements))

            _box(s, elem.get("text", ""), ex, ey, e_w, e_h,
                 fill=fill, fsz=fsz, border=border)

            elem_map[elem.get("id", f"elem_{li}_{ei}")] = {
                "x": ex, "y": ey, "w": e_w, "h": e_h
            }

    # Draw edges
    for edge in edges:
        from_id = edge.get("from", "")
        to_id = edge.get("to", "")
        f = elem_map.get(from_id)
        t = elem_map.get(to_id)
        if not f or not t:
            continue

        edge_color = _color(edge.get("color", "orange"))
        label = edge.get("label", "")
        direction = edge.get("direction", "auto")

        # Auto-detect direction
        if direction == "auto":
            dy = t["y"] - f["y"]
            dx = t["x"] - f["x"]
            if abs(dy) > abs(dx):
                direction = "down" if dy > 0 else "up"
            else:
                direction = "right" if dx > 0 else "left"

        if direction == "down":
            x1 = f["x"] + f["w"] / 2
            y1 = f["y"] + f["h"]
            x2 = t["x"] + t["w"] / 2
            y2 = t["y"]
        elif direction == "up":
            x1 = f["x"] + f["w"] / 2
            y1 = f["y"]
            x2 = t["x"] + t["w"] / 2
            y2 = t["y"] + t["h"]
        elif direction == "right":
            x1 = f["x"] + f["w"]
            y1 = f["y"] + f["h"] / 2
            x2 = t["x"]
            y2 = t["y"] + t["h"] / 2
        else:  # left
            x1 = f["x"]
            y1 = f["y"] + f["h"] / 2
            x2 = t["x"] + t["w"]
            y2 = t["y"] + t["h"] / 2

        _connector(s, x1, y1, x2, y2, color=edge_color)

        if label:
            lx = (x1 + x2) / 2 + 0.15
            ly_pos = (y1 + y2) / 2 - 0.15
            _txt(s, label, lx, ly_pos, 2.50, 0.25, sz=10,
                 color=LIGHT_GRAY, italic=True)

    # Legend
    legend = slide_spec.get("legend")
    if legend:
        _build_legend(s, legend)

    _add_slide_extras(s, slide_spec)
    return s


def _build_pipeline(prs, layout_fn, slide_spec: dict, meta: dict):
    """Pipeline: horizontal box strip at top + detail sections below."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])

    pipeline = slide_spec.get("pipeline", {})
    boxes = pipeline.get("boxes", [])
    sections = slide_spec.get("sections", [])

    # Draw pipeline strip
    if boxes:
        n = len(boxes)
        bw, bh, gap = _auto_box_size(n, 12.40)
        bh = min(bh, 0.60)  # Pipeline boxes are short
        fsz = _auto_shape_font(n)
        py = 0.88

        for i, b in enumerate(boxes):
            x = 0.40 + i * (bw + gap)
            fill = _color(b.get("fill", "dark"))
            _box(s, b.get("text", ""), x, py, bw, bh, fill=fill, fsz=fsz)
            if i < n - 1:
                _txt(s, "\u2192", x + bw, py + bh / 2 - 0.14, 0.20, 0.28,
                     sz=14, align=PP_ALIGN.CENTER)

    # Detail sections (tables or text below the pipeline)
    if sections:
        n_sec = len(sections)
        sec_gap = 0.12
        sec_w = (12.40 - (n_sec - 1) * sec_gap) / n_sec
        sec_top = 1.65

        for si, sec in enumerate(sections):
            sx = 0.20 + si * (sec_w + sec_gap)
            header = sec.get("header", "")
            _txt(s, header, sx, sec_top, sec_w, 0.30, sz=11, bold=True,
                 color=_color(sec.get("color", "white")))

            tbl = sec.get("table")
            if tbl:
                rows_data = [tbl.get("headers", [])] + tbl.get("rows", [])
                if rows_data and rows_data[0]:
                    tbl_h = min(len(rows_data) * 0.40, 4.80)
                    _c3tbl(s, sx, sec_top + 0.30, sec_w, tbl_h,
                           rows_data, fsz=9, hdr_fsz=8)

            text = sec.get("content")
            if text and not tbl:
                _txt(s, text, sx, sec_top + 0.30, sec_w, 4.50,
                     sz=_auto_body_font(text), color=WHITE)

    _add_slide_extras(s, slide_spec)
    return s


def _build_table(prs, layout_fn, slide_spec: dict, meta: dict):
    """Table slide from structured data or markdown."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    tbl = slide_spec.get("table", {})
    headers = tbl.get("headers", [])
    rows = tbl.get("rows", [])
    col_widths = tbl.get("col_widths")

    # Also handle raw markdown table in "content"
    content = slide_spec.get("content", "")
    if not headers and content and "|" in content:
        parsed = _parse_md_table(content)
        if parsed:
            headers = parsed[0]
            rows = parsed[1:]

    if not headers and not rows:
        # Fallback to text
        if content:
            _txt(s, _clean_bullets(content), 0.30, 0.95, 12.40, 5.80,
                 sz=_auto_body_font(content))
        _add_slide_extras(s, slide_spec)
        return s

    all_rows = [headers] + rows if headers else rows
    nr = len(all_rows)
    nc = len(all_rows[0]) if all_rows else 1

    # Auto column widths if not specified
    if not col_widths:
        col_widths = _auto_col_widths(all_rows, 12.40)

    tbl_top = 1.00 if not slide_spec.get("subtitle") else 1.20
    tbl_h = min(nr * 0.45, 5.80)

    # Auto font sizing for tables
    if nc >= 5 or nr >= 8:
        fsz, hdr_fsz = 9, 8
    elif nc >= 4 or nr >= 6:
        fsz, hdr_fsz = 10, 9
    else:
        fsz, hdr_fsz = 11, 10

    _c3tbl(s, 0.30, tbl_top, 12.40, tbl_h, all_rows,
           cw=col_widths, fsz=fsz, hdr_fsz=hdr_fsz)

    # Extra text below table
    extra = tbl.get("footnote") or slide_spec.get("footnote")
    if extra:
        _txt(s, extra, 0.30, tbl_top + tbl_h + 0.15, 12.40, 0.40,
             sz=10, color=LIGHT_GRAY, italic=True)

    _add_slide_extras(s, slide_spec)
    return s


def _build_comparison(prs, layout_fn, slide_spec: dict, meta: dict):
    """Side-by-side comparison (e.g., Local Dev vs Cloud)."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    sides = slide_spec.get("sides", [])
    if not sides:
        _add_slide_extras(s, slide_spec)
        return s

    n = min(len(sides), 3)
    gap = 0.50
    col_w = (12.40 - (n - 1) * gap) / n
    start_y = 1.25

    for si, side in enumerate(sides[:n]):
        sx = 0.30 + si * (col_w + gap)
        header = side.get("header", "")
        _txt(s, header, sx, start_y, col_w, 0.30, sz=14, bold=True)

        elements = side.get("elements", [])
        ey = start_y + 0.45
        for elem in elements:
            e_type = elem.get("type", "box")
            e_h = elem.get("h", 0.75)
            e_w = elem.get("w", col_w - 0.20)

            if e_type == "box":
                fill = _color(elem.get("fill", "dark"))
                _box(s, elem.get("text", ""), sx, ey, e_w, e_h,
                     fill=fill, fsz=10)
            elif e_type == "label":
                _txt(s, elem.get("text", ""), sx, ey, e_w, 0.25,
                     sz=9, color=LIGHT_GRAY, italic=True)
                e_h = 0.30
            elif e_type == "arrow":
                mid_x = sx + e_w / 2
                _connector(s, mid_x, ey, mid_x, ey + e_h,
                          color=_color(elem.get("color", "orange")))
                if elem.get("label"):
                    _txt(s, elem["label"], mid_x + 0.15, ey + 0.02,
                         1.60, 0.25, sz=9, color=LIGHT_GRAY, italic=True)

            ey += e_h + 0.08

    # Connector between sides (optional)
    bridge = slide_spec.get("bridge")
    if bridge and n >= 2:
        bx = 0.30 + col_w + gap / 2 - 0.30
        by = bridge.get("y", 3.60)
        _box(s, bridge.get("text", ""), bx, by, 1.50, 1.10,
             fill=_color(bridge.get("fill", "dark")), fsz=10)

    # Bottom table
    tbl = slide_spec.get("table")
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        all_rows = [headers] + rows
        col_widths = tbl.get("col_widths") or _auto_col_widths(all_rows, 12.73)
        tbl_top = tbl.get("top", 5.30)
        tbl_h = min(len(all_rows) * 0.40, 2.0)
        _c3tbl(s, 0.30, tbl_top, 12.73, tbl_h, all_rows,
               cw=col_widths, fsz=10, hdr_fsz=9)

    _add_slide_extras(s, slide_spec)
    return s


def _build_bullets(prs, layout_fn, slide_spec: dict, meta: dict):
    """Enhanced bullet slide with optional colored markers."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    content = slide_spec.get("content", "")
    items = slide_spec.get("items", [])

    if items:
        # Structured items with optional colors/bold
        y = 1.00
        for item in items:
            if isinstance(item, str):
                text = item
                color = WHITE
                bold = False
            else:
                text = item.get("text", "")
                color = _color(item.get("color", "white"))
                bold = item.get("bold", False)

            # Auto-prefix with bullet if not already
            if not text.startswith(("\u2022", "-", "*", "•")):
                text = "\u2022 " + text

            fsz = _auto_body_font("\n".join(
                i if isinstance(i, str) else i.get("text", "") for i in items))
            _txt(s, text, 0.50, y, 12.00, 0.40, sz=fsz, bold=bold, color=color)
            y += 0.40
    elif content:
        body = _clean_bullets(content)
        fsz = _auto_body_font(body)
        _txt(s, body, 0.30, 0.95, 12.40, 5.80, sz=fsz)

    _add_slide_extras(s, slide_spec)
    return s


def _build_text(prs, layout_fn, slide_spec: dict, meta: dict):
    """Generic text/content slide (fallback)."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    content = slide_spec.get("content", "")
    if content:
        body = _clean_bullets(content)
        fsz = _auto_body_font(body)
        _txt(s, body, 0.30, 0.95, 12.40, 5.80, sz=fsz)

    _add_slide_extras(s, slide_spec)
    return s


def _build_closing(prs, layout_fn, slide_spec: dict, meta: dict):
    """Closing / Q&A slide."""
    layouts = list(prs.slide_layouts)
    layout = layouts[14] if len(layouts) > 14 else layout_fn("Blank")
    s = prs.slides.add_slide(layout)
    title = slide_spec.get("title", "Questions?")
    content = slide_spec.get("content", "")
    text = title
    if content:
        text += "\n\n" + _clean_bullets(content)
    _txt(s, text, 0.30, 2.50, 12.40, 2.00, sz=28, bold=True)
    _add_slide_extras(s, slide_spec)
    return s


def _build_mixed(prs, layout_fn, slide_spec: dict, meta: dict):
    """Mixed content slide: renders a sequence of typed elements."""
    s = prs.slides.add_slide(layout_fn("Blank"))
    _add_title(s, slide_spec["title"])
    if slide_spec.get("subtitle"):
        _add_subtitle(s, slide_spec["subtitle"])

    elements = slide_spec.get("elements", [])
    y = 1.00

    for elem in elements:
        et = elem.get("type", "text")
        eh = elem.get("h", 1.00)

        if et == "boxes":
            # Horizontal row of boxes
            items = elem.get("items", [])
            n = len(items)
            bw, bh, gap = _auto_box_size(n, elem.get("w", 12.40))
            fsz = _auto_shape_font(n)
            for i, item in enumerate(items):
                x = elem.get("left", 0.30) + i * (bw + gap)
                fill = _color(item.get("fill", "dark"))
                _box(s, item.get("text", ""), x, y, bw, bh,
                     fill=fill, fsz=fsz)
                if i < n - 1:
                    arrow_text = elem.get("arrow", "\u2192")
                    _txt(s, arrow_text, x + bw, y + bh / 2 - 0.14,
                         0.20, 0.28, sz=14, align=PP_ALIGN.CENTER)
            eh = bh + 0.10

        elif et == "table":
            tbl = elem
            headers = tbl.get("headers", [])
            rows = tbl.get("rows", [])
            all_rows = [headers] + rows
            cw = tbl.get("col_widths") or _auto_col_widths(all_rows, elem.get("w", 12.40))
            tbl_h = min(len(all_rows) * 0.40, eh)
            _c3tbl(s, elem.get("left", 0.30), y, elem.get("w", 12.40),
                   tbl_h, all_rows, cw=cw,
                   fsz=elem.get("fsz", 10), hdr_fsz=elem.get("hdr_fsz", 9))
            eh = tbl_h + 0.10

        elif et == "text":
            _txt(s, elem.get("text", ""), elem.get("left", 0.30), y,
                 elem.get("w", 12.40), eh,
                 sz=elem.get("fsz", 12), bold=elem.get("bold", False),
                 color=_color(elem.get("color", "white")))

        elif et == "callout":
            bg = s.shapes.add_shape(1, I(elem.get("left", 0.30)), I(y),
                                     I(elem.get("w", 12.73)), I(0.52))
            bg.fill.solid()
            bg.fill.fore_color.rgb = _color(elem.get("fill", "dark"))
            bg.line.fill.background()
            label = elem.get("label", "")
            detail = elem.get("detail", "")
            if label:
                _txt(s, label, elem.get("left", 0.30) + 0.12, y + 0.04,
                     1.40, 0.24, sz=11, bold=True)
            if detail:
                _txt(s, detail, elem.get("left", 0.30) + 1.48, y + 0.04,
                     11.20, 0.44, sz=10)
            eh = 0.60

        y += eh + 0.08

    _add_slide_extras(s, slide_spec)
    return s


# ═══════════════════════════════════════════════════════════════
# Shared helpers for slide extras
# ═══════════════════════════════════════════════════════════════

def _add_slide_extras(s, slide_spec: dict):
    """Add notes, links, callout to a slide."""
    notes = slide_spec.get("notes", "")
    if notes:
        _add_notes(s, notes)

    links = slide_spec.get("links", [])
    link_y = 6.88
    for li, link in enumerate(links):
        lx = link.get("left", 0.30 + li * 4.0)
        _hlink(s, link.get("text", ""), link.get("url", ""),
               lx, link.get("top", link_y))


def _build_legend(s, legend: list, y_start: float = 6.88):
    """Color-coded legend at bottom of slide."""
    x = 1.50
    _txt(s, "Legend:", x, y_start, 1.00, 0.25, sz=10, bold=True)
    x = 2.70
    for item in legend:
        swatch = s.shapes.add_shape(1, I(x), I(y_start + 0.01), I(0.25), I(0.20))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = _color(item.get("color", "dark"))
        swatch.line.fill.background()
        _txt(s, item.get("label", ""), x + 0.32, y_start, 2.00, 0.25, sz=10)
        x += 2.60


# ═══════════════════════════════════════════════════════════════
# Utility
# ═══════════════════════════════════════════════════════════════

def _clean_bullets(text: str) -> str:
    """Convert markdown bullets to Unicode bullets."""
    lines = []
    for line in text.splitlines():
        stripped = line.strip()
        cleaned = re.sub(r"^[-*]\s+", "\u2022 ", stripped)
        # Strip bold markdown
        cleaned = re.sub(r"\*\*(.+?)\*\*", r"\1", cleaned)
        lines.append(cleaned)
    return re.sub(r"\n{3,}", "\n\n", "\n".join(lines)).strip()


def _parse_md_table(md: str) -> list[list[str]]:
    """Extract rows from markdown table."""
    rows = []
    for line in md.strip().splitlines():
        line = line.strip()
        if not line.startswith("|"):
            continue
        if re.match(r"^\|[\s\-:|]+\|$", line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        if cells:
            rows.append(cells)
    return rows


def _auto_col_widths(rows: list[list[str]], total_w: float) -> list[float]:
    """Calculate proportional column widths based on content length."""
    if not rows or not rows[0]:
        return []
    nc = len(rows[0])
    # Score each column by max cell length
    scores = []
    for ci in range(nc):
        max_len = 0
        for row in rows:
            if ci < len(row):
                max_len = max(max_len, len(str(row[ci])))
        scores.append(max(max_len, 5))  # minimum score

    total = sum(scores)
    widths = []
    for score in scores:
        w = (score / total) * total_w
        w = max(w, 0.80)  # minimum column width
        widths.append(round(w, 2))

    # Normalize to total_w
    actual_total = sum(widths)
    if actual_total > 0:
        factor = total_w / actual_total
        widths = [round(w * factor, 2) for w in widths]

    return widths


# ═══════════════════════════════════════════════════════════════
# Slide type detection (for legacy markdown input)
# ═══════════════════════════════════════════════════════════════

_CLOSING_RE = re.compile(r"question|thank\s*you|closing|q\s*&\s*a", re.IGNORECASE)


def _detect_layout(title: str, content: str, index: int, is_last: bool) -> str:
    """Auto-detect layout from markdown content (legacy mode)."""
    title_lower = title.lower()

    if index == 0 or title_lower == "title":
        return "title"
    if _CLOSING_RE.search(title_lower):
        return "closing"
    if "|" in content and re.search(r"^\|.+\|", content, re.MULTILINE):
        return "table"

    # Check for bold-labeled panels
    panels = re.split(r"(?m)^\s*\*\*(.+?)\*\*\s*$", content)
    if len(panels) >= 7:  # 3+ panels = 7+ parts
        return "panels"

    # Check for numbered steps (good candidate for flow)
    step_lines = [l for l in content.splitlines() if re.match(r"^\s*[-*]\s+\*\*(?:Step|Phase|\d)", l)]
    if len(step_lines) >= 3:
        return "flow_from_bullets"

    # Default to bullets or text
    lines = [l for l in content.splitlines() if l.strip()]
    if lines:
        bullet_lines = sum(1 for l in lines if re.match(r"^\s*[-*]\s", l))
        if bullet_lines / len(lines) >= 0.5:
            return "bullets"

    return "text"


def _legacy_panels_from_content(content: str) -> list[dict]:
    """Extract panel specs from bold-labeled markdown sections."""
    parts = re.split(r"(?m)^\s*\*\*(.+?)\*\*\s*$", content)
    panels = []
    color_map = {
        "problem": "orange", "risk": "orange", "warning": "orange",
        "constraints": "dark", "current": "dark",
        "quality": "green", "success": "green", "solution": "green", "value": "green",
    }
    for i in range(1, len(parts), 2):
        header = parts[i].strip().rstrip(":")
        body = parts[i + 1].strip() if i + 1 < len(parts) else ""
        clr = "dark"
        for kw, c in color_map.items():
            if kw in header.lower():
                clr = c
                break
        panels.append({"header": header, "color": clr, "content": _clean_bullets(body)})
    return panels


def _legacy_flow_from_bullets(content: str) -> dict:
    """Convert step-based bullet content into a flow spec."""
    items = []
    lines = content.strip().splitlines()
    for line in lines:
        stripped = line.strip()
        m = re.match(r"^[-*]\s+\*\*(.+?)\*\*[:\s]*(.*)", stripped)
        if m:
            label = m.group(1).strip()
            items.append({"text": label, "fill": "dark"})
        elif re.match(r"^[-*]\s", stripped):
            text = re.sub(r"^[-*]\s+", "", stripped)
            items.append({"text": text[:30], "fill": "dark"})

    return {
        "rows": [{"items": items[:7], "direction": "right"}],
    }


# ═══════════════════════════════════════════════════════════════
# Dispatch
# ═══════════════════════════════════════════════════════════════

_BUILDERS = {
    "title": _build_title_slide,
    "panels": _build_panels,
    "flow": _build_flow,
    "diagram": _build_diagram,
    "pipeline": _build_pipeline,
    "table": _build_table,
    "comparison": _build_comparison,
    "bullets": _build_bullets,
    "text": _build_text,
    "closing": _build_closing,
    "mixed": _build_mixed,
}


def build_deck(spec: dict[str, Any], template_path: str, output_path: str) -> str:
    """Build a C3-branded PPTX deck from a spec dict."""
    if not os.path.isfile(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    slides = spec.get("slides", [])
    if not slides:
        raise ValueError("Spec contains no slides.")

    meta = {k: v for k, v in spec.items() if k != "slides"}
    prs, layout_fn = _load_template(template_path)

    for si, slide_spec in enumerate(slides):
        layout = slide_spec.get("layout", "")
        title = slide_spec.get("title", "")

        # Auto-detect layout for legacy markdown slides
        if not layout:
            content = slide_spec.get("content", "")
            is_last = si == len(slides) - 1
            layout = _detect_layout(title, content, si, is_last)

            # Convert detected layout to structured spec
            if layout == "panels":
                slide_spec["panels"] = _legacy_panels_from_content(content)
            elif layout == "flow_from_bullets":
                flow_spec = _legacy_flow_from_bullets(content)
                slide_spec.update(flow_spec)
                layout = "flow"
            elif layout == "table":
                # Parse markdown table into structured form
                parsed = _parse_md_table(content)
                if parsed:
                    slide_spec["table"] = {
                        "headers": parsed[0],
                        "rows": parsed[1:],
                    }
                    # Extra non-table text becomes notes
                    non_table = []
                    in_table = False
                    for line in content.splitlines():
                        if line.strip().startswith("|"):
                            in_table = True
                            continue
                        if in_table and not line.strip():
                            in_table = False
                        if not in_table and line.strip():
                            non_table.append(line.strip())
                    if non_table and not slide_spec.get("notes"):
                        slide_spec["notes"] = "\n".join(non_table)

        builder = _BUILDERS.get(layout, _build_text)
        builder(prs, layout_fn, slide_spec, meta)

    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)

    prs.save(output_path)
    return os.path.abspath(output_path)


# ═══════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="C3 PPT Adapter v2 — rich content engine.")
    parser.add_argument("input", help="Path to YAML or markdown file.")
    parser.add_argument("--template",
                        default=os.path.join(os.path.dirname(__file__),
                                             "templates", "c3-template.pptx"),
                        help="C3 PPTX template path.")
    parser.add_argument("--output", default=None, help="Output PPTX path.")
    args = parser.parse_args()

    input_path = os.path.abspath(args.input)
    if not os.path.isfile(input_path):
        print(f"Error: file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    spec = _parse_input(input_path)
    slide_count = len(spec.get("slides", []))
    title = spec.get("title", "(untitled)")
    print(f'Parsed "{title}" — {slide_count} slides.')

    if slide_count == 0:
        print("Error: no slides found.", file=sys.stderr)
        sys.exit(1)

    if args.output:
        output_path = os.path.abspath(args.output)
    else:
        stem = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(os.path.dirname(input_path),
                                   f"{stem}_deck.pptx")

    template_path = os.path.abspath(args.template)
    try:
        saved = build_deck(spec, template_path, output_path)
        print(f"\u2713 Saved: {saved}")
    except Exception as exc:
        print(f"Error: {exc}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
