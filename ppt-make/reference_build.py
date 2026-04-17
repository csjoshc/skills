#!/usr/bin/env python3
"""
C3 PPT Reference Build — semi-modular.

Structure:
  1. C3Palette       — colour constants
  2. Sources         — URLs for speaker notes & hyperlinks
  3. C3Helpers       — stateless helper functions (importable)
  4. load_clean_template — template loader with footer-year fix
  5. Slide builders   — one function per slide (composable)
  6. main()          — orchestrates template load + slide assembly

To reuse in a new deck:
    from reference_build import C3Palette, C3Helpers, Sources, load_clean_template
"""
from pptx import Presentation
from pptx.util import Inches as I, Pt as P
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE


# ═══════════════════════════════════════════════════════════════
# SECTION 1 — Palette
# ═══════════════════════════════════════════════════════════════

class C3Palette:
    """C3 dark-theme colour constants."""
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xC8, 0xC8, 0xC8)
    MED_GRAY   = RGBColor(0x8C, 0x8C, 0x8C)
    DARK_FILL  = RGBColor(0x50, 0x50, 0x50)
    GRAY_FILL  = RGBColor(0x8C, 0x8C, 0x8C)
    C3_BLUE    = RGBColor(0x06, 0xA7, 0xE0)
    ORANGE     = RGBColor(0xF7, 0x94, 0x30)
    GREEN      = RGBColor(0x9C, 0xCD, 0x6C)
    WARM_GRAY  = RGBColor(0xCB, 0xC8, 0xC7)
    NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
    TBL_HDR    = RGBColor(0x50, 0x50, 0x50)
    TBL_ROW_A  = RGBColor(0x3A, 0x3A, 0x3A)
    TBL_ROW_B  = RGBColor(0x2D, 0x2D, 0x2D)
    PANEL_DARK = RGBColor(0x35, 0x35, 0x35)

P_ = C3Palette  # short alias


# ═══════════════════════════════════════════════════════════════
# SECTION 2 — Source URLs
# ═══════════════════════════════════════════════════════════════

class Sources:
    """Project-specific source URLs.

    Override this class or pass a sources dict via the adapter for
    project-specific links.  The default values are placeholders from
    the original integration-demo deck — replace per project.
    """

    @classmethod
    def from_dict(cls, mapping: dict) -> "Sources":
        """Create a Sources instance from a flat dict of label→url pairs.

        Useful when sources come from YAML frontmatter rather than
        being hardcoded::

            srcs = Sources.from_dict({"wiki": "https://...", "jira": "https://..."})
            srcs.get("wiki")  # → "https://..."
        """
        instance = cls()
        instance._custom = mapping
        return instance

    def __init__(self):
        self._custom: dict = {}

    def get(self, key: str, default: str = "") -> str:
        """Look up a source URL by key (case-insensitive)."""
        key_lower = key.lower()
        # Check custom sources first
        for k, v in self._custom.items():
            if k.lower() == key_lower:
                return v
        # Fall back to class-level attributes
        for attr in dir(self):
            if attr.lower() == key_lower and not attr.startswith("_"):
                return getattr(self, attr, default)
        return default

    # --- Default project sources (override per project) ---
    PLACEHOLDER = "(add project-specific sources via Sources.from_dict or YAML frontmatter)"

S = Sources  # short alias


# ═══════════════════════════════════════════════════════════════
# SECTION 3 — Helpers
# ═══════════════════════════════════════════════════════════════

class C3Helpers:
    """Stateless helpers for building C3-branded slides."""

    @staticmethod
    def add_title(s, text, sz=24):
        tb = s.shapes.add_textbox(I(0.30), I(0.18), I(12.75), I(0.58))
        tf = tb.text_frame; tf.word_wrap = False
        r = tf.paragraphs[0].add_run(); r.text = text
        r.font.size = P(sz); r.font.bold = True; r.font.color.rgb = P_.WHITE

    @staticmethod
    def txt(s, text, left, top, w, h, sz=12, bold=False, color=None,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
        color = color or P_.WHITE
        tb = s.shapes.add_textbox(I(left), I(top), I(w), I(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        for i, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run(); r.text = line
            r.font.size = P(sz); r.font.bold = bold
            r.font.italic = italic; r.font.color.rgb = color

    @staticmethod
    def box(s, text, left, top, w, h, fill=None, fsz=12, fbold=True,
            fcolor=None, border=None):
        fcolor = fcolor or P_.WHITE
        sh = s.shapes.add_shape(1, I(left), I(top), I(w), I(h))
        if fill:
            sh.fill.solid(); sh.fill.fore_color.rgb = fill
        else:
            sh.fill.background()
        if border:
            sh.line.color.rgb = border; sh.line.width = P(0.75)
        else:
            sh.line.fill.background()
        if text:
            tf = sh.text_frame; tf.word_wrap = True
            for j, line in enumerate(text.split('\n')):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = line
                r.font.size = P(fsz); r.font.bold = fbold; r.font.color.rgb = fcolor
        return sh

    @staticmethod
    def c3tbl(s, left, top, w, h, rows, cw=None, fsz=11, hdr_fsz=None):
        """C3-styled table. hdr_fsz controls header font (defaults to fsz)."""
        hdr_fsz = hdr_fsz or fsz
        nr, nc = len(rows), len(rows[0])
        t = s.shapes.add_table(nr, nc, I(left), I(top), I(w), I(h)).table
        if cw:
            for i, c in enumerate(cw): t.columns[i].width = I(c)
        for ri, row in enumerate(rows):
            for ci, cell_text in enumerate(row):
                cell = t.cell(ri, ci); cell.text = str(cell_text)
                p = cell.text_frame.paragraphs[0]
                r = p.runs[0] if p.runs else p.add_run()
                r.font.color.rgb = P_.WHITE
                if ri == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = P_.TBL_HDR
                    r.font.bold = True; r.font.size = P(hdr_fsz)
                else:
                    r.font.size = P(fsz)
                    if ri % 2 == 1:
                        cell.fill.solid(); cell.fill.fore_color.rgb = P_.TBL_ROW_A
                    else:
                        cell.fill.solid(); cell.fill.fore_color.rgb = P_.TBL_ROW_B
        return t

    @staticmethod
    def add_notes(s, text):
        tf = s.notes_slide.notes_text_frame; tf.text = text

    @staticmethod
    def hlink_txt(s, text, url, left, top, w, h, sz=8, bold=False, color=None,
                  align=PP_ALIGN.LEFT):
        """Clickable hyperlink textbox (defaults to C3 Blue, 8pt)."""
        color = color or P_.C3_BLUE
        tb = s.shapes.add_textbox(I(left), I(top), I(w), I(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = P(sz); r.font.bold = bold; r.font.color.rgb = color
        r.hyperlink.address = url
        return tb

    @staticmethod
    def connector(s, x1, y1, x2, y2, color=None, width=1.5):
        """Straight connector with directional arrowhead at end."""
        from pptx.oxml.ns import qn
        from lxml import etree
        color = color or P_.ORANGE
        conn = s.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
        conn.line.color.rgb = color; conn.line.width = P(width)
        ln = conn.line._ln
        tail = ln.find(qn('a:tailEnd'))
        if tail is None:
            tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
        return conn

    @staticmethod
    def arrow_down(s, x, y, sz=14):
        C3Helpers.txt(s, "↓", x, y, 0.30, 0.30, sz=sz, color=P_.ORANGE)

    @staticmethod
    def panel_column(s, x, w, hdr, hdr_color, content, body_sz=14):
        """Dark panel with coloured header bar (common 3-column pattern)."""
        bg = s.shapes.add_shape(1, I(x), I(0.92), I(w), I(5.80))
        bg.fill.solid(); bg.fill.fore_color.rgb = P_.PANEL_DARK
        bg.line.color.rgb = P_.WARM_GRAY; bg.line.width = P(0.5)
        hb = s.shapes.add_shape(1, I(x), I(0.92), I(w), I(0.42))
        hb.fill.solid(); hb.fill.fore_color.rgb = hdr_color; hb.line.fill.background()
        C3Helpers.txt(s, hdr, x+0.12, 0.95, w-0.24, 0.36, sz=14, bold=True, color=P_.WHITE)
        C3Helpers.txt(s, content, x+0.15, 1.50, w-0.30, 5.00, sz=body_sz, wrap=True, color=P_.WHITE)

h = C3Helpers  # short alias


# ═══════════════════════════════════════════════════════════════
# SECTION 4 — Template management
# ═══════════════════════════════════════════════════════════════

def load_clean_template(path):
    """Load clean C3 template, fix footer year, return (prs, layout_fn)."""
    prs = Presentation(path)
    import datetime
    year = str(datetime.date.today().year)
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '2025' in run.text:
                            run.text = run.text.replace('2025', year)
        for lay in master.slide_layouts:
            for shape in lay.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if '2025' in run.text:
                                run.text = run.text.replace('2025', year)

    def layout(name):
        for L in prs.slide_layouts:
            if L.name == name: return L
        return prs.slide_layouts[13]

    return prs, layout


# ═══════════════════════════════════════════════════════════════
# SECTION 5 — Slide builders  (each returns the slide object)
# ═══════════════════════════════════════════════════════════════
# Only the first three are shown fully as examples.
# The remaining slides follow the same pattern — see main().

# ═══════════════════════════════════════════════════════════════
# SECTION 5 — Example slide builders
# ═══════════════════════════════════════════════════════════════
# These demonstrate the patterns. For automated slide generation
# from reframe markdown, use adapter.py instead.

def example_title_slide(prs, layout_fn, title="Untitled Deck",
                        subtitle="", author_line=""):
    """Example: build a title slide from parameters (not hardcoded)."""
    s = prs.slides.add_slide(layout_fn("Title Slide Option"))
    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
            if ph.text_frame.paragraphs[0].runs:
                ph.text_frame.paragraphs[0].runs[0].font.size = P(36)
        elif idx == 15:
            ph.text = subtitle
            if ph.text_frame.paragraphs[0].runs:
                ph.text_frame.paragraphs[0].runs[0].font.size = P(18)
        elif idx == 16:
            ph.text = author_line
            if ph.text_frame.paragraphs[0].runs:
                ph.text_frame.paragraphs[0].runs[0].font.size = P(14)
    return s


def example_panel_slide(prs, layout_fn, title, panels):
    """Example: 3-column panel slide.

    Args:
        panels: list of (header, color, content) tuples.
    """
    s = prs.slides.add_slide(layout_fn("Blank"))
    h.add_title(s, title)
    CW, CG = 4.10, 0.12
    for i, (hdr, clr, content) in enumerate(panels[:3]):
        h.panel_column(s, 0.30 + i * (CW + CG), CW, hdr, clr, content)
    return s


def example_diagram_slide(prs, layout_fn, title, subtitle_text, shapes):
    """Example: architecture diagram built from shape definitions.

    Args:
        shapes: list of dicts with keys: text, left, top, w, h, fill, [fsz, border]
    """
    s = prs.slides.add_slide(layout_fn("Blank"))
    h.add_title(s, title)
    if subtitle_text:
        h.txt(s, subtitle_text, 0.30, 0.82, 12.0, 0.30, sz=12,
              color=P_.LIGHT_GRAY, italic=True)
    for sh in shapes:
        h.box(s, sh["text"], sh["left"], sh["top"], sh["w"], sh["h"],
              fill=sh.get("fill"), fsz=sh.get("fsz", 12),
              border=sh.get("border"))
    return s


# ═══════════════════════════════════════════════════════════════
# SECTION 6 — Main
# ═══════════════════════════════════════════════════════════════
# This file is the importable foundation for C3 PPTX generation.
# For automated markdown-to-deck conversion, use adapter.py.

if __name__ == "__main__":
    print("C3 PPT Reference Module")
    print("Importable: C3Palette, C3Helpers, Sources, load_clean_template")
    print("For markdown→PPTX: python adapter.py input.md")
